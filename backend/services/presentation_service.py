import subprocess
from pathlib import Path

from utils.file_utils import temp_path
from utils.platform_utils import get_libreoffice_path
from utils.zip_utils import bundle_to_zip


def _lo_convert(inp: Path, out_dir: Path, target_fmt: str) -> Path:
    soffice = get_libreoffice_path()
    result = subprocess.run(
        [soffice, "--headless", "--convert-to", target_fmt, "--outdir", str(out_dir), str(inp)],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice error: {result.stderr[-500:]}")
    out = out_dir / (inp.stem + f".{target_fmt}")
    if not out.exists():
        candidates = list(out_dir.glob(f"*.{target_fmt}"))
        if not candidates:
            raise RuntimeError(f"LibreOffice produced no .{target_fmt} output")
        out = candidates[0]
    return out


def pptx_to_pdf(inp: Path, out: Path) -> None:
    tmp_dir = inp.parent
    pdf = _lo_convert(inp, tmp_dir, "pdf")
    if pdf != out:
        pdf.rename(out)


def pptx_to_images_zip(inp: Path, out_zip: Path) -> None:
    """Export each slide as PNG via LibreOffice impress export."""
    tmp_dir = temp_path("").parent / f"slides_{inp.stem}"
    tmp_dir.mkdir(exist_ok=True)
    try:
        soffice = get_libreoffice_path()
        subprocess.run(
            [soffice, "--headless", "--convert-to", "png", "--outdir", str(tmp_dir), str(inp)],
            capture_output=True, text=True, timeout=120, check=True,
        )
        slides = sorted(tmp_dir.glob("*.png"))
        bundle_to_zip([(p, p.name) for p in slides], out_zip)
    finally:
        for f in tmp_dir.glob("*"):
            f.unlink(missing_ok=True)
        tmp_dir.rmdir()


def pptx_to_html(inp: Path, out: Path) -> None:
    """Extract slide text as basic HTML via python-pptx."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx")
    prs = Presentation(str(inp))
    parts = ['<!DOCTYPE html><html><head><meta charset="utf-8"></head><body>']
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"<section><h2>Slide {i}</h2>")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(f"<p>{text}</p>")
        parts.append("</section><hr>")
    parts.append("</body></html>")
    out.write_text("\n".join(parts), encoding="utf-8")


def office_to_pptx(inp: Path, out: Path) -> None:
    """Convert PPT, ODP, or PDF to PPTX via LibreOffice."""
    tmp_dir = inp.parent
    pptx = _lo_convert(inp, tmp_dir, "pptx")
    if pptx != out:
        pptx.rename(out)


def office_to_pdf(inp: Path, out: Path) -> None:
    tmp_dir = inp.parent
    pdf = _lo_convert(inp, tmp_dir, "pdf")
    if pdf != out:
        pdf.rename(out)


def compress_pptx(inp: Path, out: Path, max_image_px: int = 1280) -> None:
    """Recompress embedded images in a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image  # type: ignore[import-untyped]
        import io
    except ImportError:
        raise RuntimeError("python-pptx and Pillow are required for PPTX compression.")

    prs = Presentation(str(inp))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                continue
            try:
                img_blob = shape.image.blob
                img = Image.open(io.BytesIO(img_blob))
                if max(img.size) > max_image_px:
                    img.thumbnail((max_image_px, max_image_px), Image.LANCZOS)
                    buf = io.BytesIO()
                    fmt = img.format or "JPEG"
                    if img.mode == "RGBA":
                        img = img.convert("RGB")
                        fmt = "JPEG"
                    img.save(buf, format=fmt, quality=80, optimize=True)
                    shape.image._blob = buf.getvalue()
            except Exception:
                continue  # skip shapes that can't be recompressed
    prs.save(str(out))


def merge_pptx(input_paths: list[Path], out: Path) -> None:
    """Merge multiple PPTX files into one."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        import copy
    except ImportError:
        raise RuntimeError("python-pptx is required.")

    base_prs = Presentation(str(input_paths[0]))

    for extra_path in input_paths[1:]:
        extra_prs = Presentation(str(extra_path))
        for slide in extra_prs.slides:
            slide_layout = base_prs.slide_layouts[0]
            new_slide = base_prs.slides.add_slide(slide_layout)
            # Clear default placeholders from layout
            for ph in new_slide.placeholders:
                sp = ph._element
                sp.getparent().remove(sp)
            # Deep-copy shapes from source slide
            for shape in slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.insert(2, el)

    base_prs.save(str(out))
